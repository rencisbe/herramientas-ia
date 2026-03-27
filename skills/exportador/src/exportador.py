import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


def agregar_grafica_cpu(
    df: pd.DataFrame,
    template_path: str,
    output_path: str,
    slide_index: int = 1,
) -> None:
    prs = Presentation(template_path)
    slide = prs.slides[slide_index]

    chart_data = CategoryChartData()
    chart_data.categories = df["TheDate"].astype(str).tolist()
    chart_data.add_series("CPU %", tuple(df["CPUUExec"].tolist()))

    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4)
    slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data)

    prs.save(output_path)

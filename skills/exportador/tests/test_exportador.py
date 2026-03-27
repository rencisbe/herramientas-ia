import pytest
import pandas as pd
import os
from pptx import Presentation
from src.exportador import agregar_grafica_cpu

def test_agregar_grafica_cpu_crea_archivo_valido(tmp_path):
    # Datos simulados
    df_cpu = pd.DataFrame({
        "TheDate": ["2026-03-25", "2026-03-26", "2026-03-27"],
        "CPUUExec": [45.2, 88.5, 60.1]
    })
    
    # Rutas adaptadas a tu nueva estructura
    template_path = "templates/Monthly_Report_Template.pptx"
    output_path = tmp_path / "Reporte_Salida_Test.pptx"
    
    assert os.path.exists(template_path), "Falta el archivo PPTX en la carpeta templates/"

    # Ejecutar la función
    agregar_grafica_cpu(
        df=df_cpu,
        template_path=template_path,
        output_path=str(output_path),
        slide_index=1
    )
    
    # Verificaciones
    assert os.path.exists(output_path), "No se generó el archivo de salida"
    prs = Presentation(output_path)
    assert len(prs.slides) > 1, "La plantilla no tiene suficientes slides"
    
    slide = prs.slides[1]
    tiene_grafico = any(shape.has_chart for shape in slide.shapes)
    assert tiene_grafico, "No se insertó el gráfico de CPU"